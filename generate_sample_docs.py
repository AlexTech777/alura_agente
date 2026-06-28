"""
Genera los 8 documentos de ejemplo de NimbusFlow (plataforma SaaS ficticia),
uno por cada formato soportado por el agente: PDF, Word, Excel, PowerPoint,
Markdown, CSV, JSON y HTML.

Ejecutar con: python generate_sample_docs.py
"""
import os
import json

os.makedirs("data", exist_ok=True)

# ----------------------------------------------------------------------------
# 1. PDF - Base de conocimiento del producto
# ----------------------------------------------------------------------------
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet

styles = getSampleStyleSheet()
doc = SimpleDocTemplate("data/base_conocimiento_producto.pdf", pagesize=letter)
story = [Paragraph("NimbusFlow - Base de Conocimiento del Producto", styles["Title"]), Spacer(1, 14)]
secciones_pdf = [
    ("Que es NimbusFlow",
     "NimbusFlow es una plataforma SaaS de gestion de proyectos y automatizacion de "
     "tareas para equipos de hasta 500 personas. Permite crear tableros Kanban, "
     "diagramas de Gantt, automatizaciones sin codigo y reportes de productividad en "
     "tiempo real."),
    ("Modulos principales",
     "La plataforma cuenta con cuatro modulos: Tableros (gestion visual de tareas), "
     "Automatizaciones (flujos disparados por eventos), Reportes (paneles de metricas "
     "de equipo) e Integraciones (conexion con Slack, Google Drive, GitHub y Microsoft Teams)."),
    ("Roles de usuario",
     "Existen tres roles dentro de un espacio de trabajo: Administrador (gestiona "
     "facturacion, usuarios y permisos), Editor (puede crear y modificar tareas y "
     "tableros) e Invitado (solo puede ver y comentar, sin permisos de edicion)."),
    ("Limites tecnicos",
     "El plan Free permite hasta 3 tableros y 10 automatizaciones activas. El tamano "
     "maximo por archivo adjunto es de 250 MB en el plan Pro y de 1 GB en el plan Business."),
    ("Disponibilidad del servicio",
     "NimbusFlow garantiza un SLA de disponibilidad del 99.9% mensual para los planes "
     "Business. El estado del servicio puede consultarse en status.nimbusflow.com."),
    ("Exportacion de datos",
     "Los usuarios pueden exportar sus tableros en formato CSV o JSON en cualquier "
     "momento desde Configuracion > Exportar datos, sin costo adicional."),
]
for h, b in secciones_pdf:
    story.append(Paragraph(h, styles["Heading2"]))
    story.append(Spacer(1, 6))
    story.append(Paragraph(b, styles["Normal"]))
    story.append(Spacer(1, 12))
doc.build(story)
print("Generado: data/base_conocimiento_producto.pdf")

# ----------------------------------------------------------------------------
# 2. Word - FAQ de soporte
# ----------------------------------------------------------------------------
from docx import Document as DocxDocument

docx_doc = DocxDocument()
docx_doc.add_heading("NimbusFlow - Preguntas Frecuentes de Soporte", level=0)
faqs = [
    ("Como recupero mi contrasena",
     "Ingresa a la pantalla de inicio de sesion y hace clic en 'Olvide mi contrasena'. "
     "Recibiras un correo con un enlace valido por 30 minutos para crear una nueva clave."),
    ("Como invito a un nuevo miembro al equipo",
     "Desde Configuracion > Miembros, hace clic en 'Invitar persona', ingresa su correo "
     "electronico y seleccionas el rol. El invitado recibe un enlace valido por 7 dias."),
    ("Por que no puedo crear mas tableros",
     "Esto ocurre si tu espacio de trabajo esta en el plan Free, que tiene un limite de "
     "3 tableros activos. Para crear mas, actualiza a un plan Pro o Business."),
    ("Como cancelo una automatizacion",
     "Ve al modulo Automatizaciones, selecciona la automatizacion y hace clic en "
     "'Pausar' o 'Eliminar'."),
    ("Cuanto tiempo tarda en responder soporte",
     "4 horas habiles para usuarios Pro, y 1 hora habil para usuarios Business "
     "(soporte prioritario)."),
    ("Puedo migrar datos desde otra herramienta",
     "Si, hay importadores nativos desde Trello, Asana y hojas de calculo CSV."),
    ("Como contacto a soporte",
     "Por soporte@nimbusflow.com o chat en vivo, lunes a viernes de 9 a 18 hs (UTC-3)."),
]
for pregunta, respuesta in faqs:
    docx_doc.add_heading(pregunta, level=2)
    docx_doc.add_paragraph(respuesta)
docx_doc.save("data/faq_soporte.docx")
print("Generado: data/faq_soporte.docx")

# ----------------------------------------------------------------------------
# 3. Excel - Planes y precios
# ----------------------------------------------------------------------------
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment

wb = Workbook()
ws = wb.active
ws.title = "Planes y Precios"
headers = ["Plan", "Precio mensual (USD/usuario)", "Precio anual (USD/usuario/mes)",
           "Tableros", "Automatizaciones", "Almacenamiento", "Soporte", "SLA disponibilidad"]
ws.append(headers)
for cell in ws[1]:
    cell.font = Font(bold=True)
    cell.alignment = Alignment(horizontal="center")
rows = [
    ["Free", 0, 0, "Hasta 3", "Hasta 10", "1 GB total", "Correo (48 hs)", "No aplica"],
    ["Pro", 9, 7, "Ilimitados", "Ilimitadas", "50 GB por espacio", "Chat (4 hs habiles)", "No aplica"],
    ["Business", 18, 15, "Ilimitados", "Ilimitadas", "1 TB por espacio", "Prioritario (1 hora habil)", "99.9% mensual"],
    ["Enterprise", "Personalizado", "Personalizado", "Ilimitados", "Ilimitadas",
     "Personalizado", "Gerente de cuenta dedicado", "SLA personalizado"],
]
for row in rows:
    ws.append(row)
for col in ws.columns:
    max_len = max(len(str(c.value)) for c in col if c.value is not None)
    ws.column_dimensions[col[0].column_letter].width = max_len + 4

ws2 = wb.create_sheet("Notas adicionales")
notas = [
    ["Concepto", "Detalle"],
    ["Periodo de prueba", "14 dias gratis en planes Pro y Business, sin tarjeta de credito"],
    ["Descuento ONG / educacion", "30% sobre planes Pro y Business, sujeto a verificacion"],
    ["Plan al vencer la prueba", "Vuelve automaticamente a Free si no se activa una suscripcion"],
]
for row in notas:
    ws2.append(row)
for cell in ws2[1]:
    cell.font = Font(bold=True)
ws2.column_dimensions["A"].width = 28
ws2.column_dimensions["B"].width = 60
wb.save("data/planes_y_precios.xlsx")
print("Generado: data/planes_y_precios.xlsx")

# ----------------------------------------------------------------------------
# 4. PowerPoint - Roadmap estrategico
# ----------------------------------------------------------------------------
from pptx import Presentation

prs = Presentation()

def add_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullets[0]
    for b in bullets[1:]:
        p = body.add_paragraph()
        p.text = b

add_slide("NimbusFlow - Roadmap Estrategico 2026",
           ["Vision: ser la plataforma de gestion de proyectos preferida para equipos remotos de Latinoamerica"])
add_slide("Q1 2026 - Fundaciones", [
    "Lanzamiento de SSO para el plan Business",
    "Mejora del motor de automatizaciones: disparadores por fecha",
    "Integracion nativa con Google Calendar",
])
add_slide("Q2 2026 - Inteligencia Artificial", [
    "Lanzamiento de NimbusFlow AI: resumen automatico de tableros",
    "Sugerencias automaticas de fechas limite",
    "Asistente conversacional dentro de la plataforma",
])
add_slide("Q3 2026 - Expansion regional", [
    "Soporte multi-idioma: portugues e ingles",
    "Alianzas comerciales en Brasil y Mexico",
    "Plan Enterprise con despliegue on-premise",
])
add_slide("Q4 2026 - Escalabilidad", [
    "Migracion de infraestructura critica a OCI",
    "Meta: SLA de disponibilidad del 99.95%",
    "Certificacion ISO 27001",
])
add_slide("Metricas objetivo para fin de 2026", [
    "50.000 espacios de trabajo activos",
    "Retencion anual superior al 90% en plan Business",
    "Soporte prioritario con respuesta menor a 30 minutos",
])
prs.save("data/roadmap_producto.pptx")
print("Generado: data/roadmap_producto.pptx")

# ----------------------------------------------------------------------------
# 5. Markdown - Terminos de uso
# ----------------------------------------------------------------------------
md_content = """# NimbusFlow - Términos y Condiciones de Uso

## 1. Aceptación de los términos
Al crear una cuenta en NimbusFlow, la persona usuaria acepta estos términos de uso.

## 2. Uso permitido
La plataforma puede utilizarse para gestionar proyectos y flujos de trabajo internos.
Queda prohibido su uso para contenido ilegal o que infrinja derechos de terceros.

## 3. Cancelación de la cuenta
Se puede cancelar en cualquier momento desde Configuración > Facturación. La cancelación
tiene efecto al finalizar el período ya pagado.

## 4. Política de reembolsos
Reembolso completo si la solicitud se hace dentro de los primeros 14 días desde el primer pago.

## 5. Suspensión del servicio
NimbusFlow puede suspender una cuenta ante un uso que viole estos términos, notificando
dentro de las 24 horas.

## 6. Propiedad de los datos
Todo el contenido cargado por la persona usuaria le pertenece exclusivamente a ella.

## 7. Modificaciones de los términos
Los cambios sustanciales se notifican con al menos 15 días de anticipación.
"""
with open("data/terminos_de_uso.md", "w", encoding="utf-8") as f:
    f.write(md_content)
print("Generado: data/terminos_de_uso.md")

# ----------------------------------------------------------------------------
# 6. CSV - Tickets de soporte
# ----------------------------------------------------------------------------
csv_content = """id_ticket,fecha,plan_cliente,categoria,prioridad,tiempo_resolucion_horas,estado
T-1001,2026-01-12,Pro,Facturacion,Media,3.5,Resuelto
T-1002,2026-01-15,Business,Bug,Alta,0.8,Resuelto
T-1003,2026-01-20,Free,Consulta general,Baja,40.0,Resuelto
T-1004,2026-02-02,Pro,Integraciones,Media,5.2,Resuelto
T-1005,2026-02-10,Business,Bug,Alta,1.1,Resuelto
T-1006,2026-02-18,Pro,Automatizaciones,Media,6.0,Resuelto
T-1007,2026-03-01,Free,Facturacion,Baja,45.3,Resuelto
T-1008,2026-03-05,Business,Rendimiento,Alta,0.9,Resuelto
T-1009,2026-03-14,Pro,Consulta general,Media,4.7,Resuelto
T-1010,2026-03-22,Business,Seguridad,Alta,0.5,Resuelto
T-1011,2026-04-02,Free,Consulta general,Baja,38.9,Resuelto
T-1012,2026-04-10,Pro,Integraciones,Media,5.8,En progreso
T-1013,2026-04-15,Business,Bug,Alta,1.3,Resuelto
T-1014,2026-04-28,Pro,Automatizaciones,Media,4.9,Resuelto
T-1015,2026-05-05,Business,Facturacion,Media,2.0,Resuelto
"""
with open("data/tickets_soporte.csv", "w", encoding="utf-8") as f:
    f.write(csv_content)
print("Generado: data/tickets_soporte.csv")

# ----------------------------------------------------------------------------
# 7. JSON - Especificacion de API de integraciones
# ----------------------------------------------------------------------------
api_spec = {
    "nombre_api": "NimbusFlow Integrations API",
    "version": "2.3.0",
    "base_url": "https://api.nimbusflow.com/v2",
    "autenticacion": {
        "tipo": "API Key",
        "header": "X-NimbusFlow-Key",
        "nota": "Solo el rol Administrador puede generar claves, desde Configuracion > Desarrolladores."
    },
    "limites_de_uso": {
        "plan_pro": "1000 solicitudes por hora",
        "plan_business": "10000 solicitudes por hora",
        "plan_enterprise": "Limite personalizado segun acuerdo"
    },
    "endpoints": [
        {"metodo": "GET", "ruta": "/tableros", "descripcion": "Lista los tableros del espacio de trabajo."},
        {"metodo": "POST", "ruta": "/tableros/{id}/tareas", "descripcion": "Crea una tarea en un tablero."},
        {"metodo": "GET", "ruta": "/automatizaciones", "descripcion": "Lista las automatizaciones activas y pausadas."},
        {"metodo": "POST", "ruta": "/webhooks", "descripcion": "Registra un webhook de eventos en tiempo real."}
    ],
    "integraciones_nativas_disponibles": ["Slack", "Google Drive", "GitHub", "Microsoft Teams"],
    "soporte_tecnico_api": {
        "contacto": "developers@nimbusflow.com",
        "documentacion_completa": "https://developers.nimbusflow.com/docs"
    }
}
with open("data/api_integraciones.json", "w", encoding="utf-8") as f:
    json.dump(api_spec, f, ensure_ascii=False, indent=2)
print("Generado: data/api_integraciones.json")

# ----------------------------------------------------------------------------
# 8. HTML - Politica de privacidad
# ----------------------------------------------------------------------------
html_content = """<!DOCTYPE html>
<html lang="es">
<head><meta charset="UTF-8"><title>NimbusFlow - Política de Privacidad</title></head>
<body>
<h1>NimbusFlow - Política de Privacidad</h1>
<h2>Datos que recopilamos</h2>
<p>Recopilamos datos de cuenta, datos de uso y datos tecnicos (IP, navegador, dispositivo).</p>
<h2>Finalidad del tratamiento de datos</h2>
<p>Brindar el servicio, mejorar la plataforma y, con consentimiento, enviar marketing.</p>
<h2>Conservacion de datos</h2>
<p>Los datos se conservan mientras la cuenta este activa, y hasta 90 dias tras su eliminacion.</p>
<h2>Compartir datos con terceros</h2>
<p>No se venden datos a terceros. Se comparten solo con proveedores de infraestructura.</p>
<h2>Derechos del usuario</h2>
<p>Se puede solicitar acceso, rectificacion o eliminacion escribiendo a privacidad@nimbusflow.com.</p>
<h2>Seguridad de la informacion</h2>
<p>Toda la informacion se transmite cifrada con TLS 1.2+ y se almacena cifrada en reposo.</p>
</body>
</html>
"""
with open("data/politica_privacidad.html", "w", encoding="utf-8") as f:
    f.write(html_content)
print("Generado: data/politica_privacidad.html")

print("\nListo. Se generaron los 8 documentos de ejemplo en data/.")
