"""
loaders.py - Lectores de documentos para el Alura Agente.

Cada funcion `load_*` recibe la ruta de un archivo y devuelve una lista de
langchain_core.documents.Document, con el texto extraido y metadata indicando
el archivo de origen. Esto permite que agent.py trate todos los formatos de
manera uniforme una vez extraido el texto.

Formatos soportados: PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx),
Markdown (.md), CSV (.csv), JSON (.json) y HTML (.html/.htm).
"""
import os
import csv
import json
import glob

from langchain_core.documents import Document
from pypdf import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup


def load_pdf(path: str):
    reader = PdfReader(path)
    docs = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            docs.append(Document(page_content=text, metadata={"fuente": os.path.basename(path), "pagina": i + 1}))
    return docs


def load_docx(path: str):
    doc = DocxDocument(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    text = "\n".join(parts)
    return [Document(page_content=text, metadata={"fuente": os.path.basename(path)})]


def load_xlsx(path: str):
    wb = load_workbook(path, data_only=True)
    docs = []
    for sheet in wb.worksheets:
        lines = [f"Hoja: {sheet.title}"]
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                lines.append(" | ".join(str(c) if c is not None else "" for c in row))
        text = "\n".join(lines)
        docs.append(Document(page_content=text, metadata={"fuente": os.path.basename(path), "hoja": sheet.title}))
    return docs


def load_pptx(path: str):
    prs = Presentation(path)
    docs = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        texts.append(line)
        text = "\n".join(texts)
        if text.strip():
            docs.append(Document(page_content=text, metadata={"fuente": os.path.basename(path), "diapositiva": i + 1}))
    return docs


def load_markdown(path: str):
    with open(path, "r", encoding="utf-8") as f:
        text = f.read()
    return [Document(page_content=text, metadata={"fuente": os.path.basename(path)})]


def load_csv(path: str):
    with open(path, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        rows = list(reader)
    header = ", ".join(reader.fieldnames or [])
    lines = [f"Columnas: {header}"]
    for row in rows:
        lines.append(", ".join(f"{k}: {v}" for k, v in row.items()))
    text = "\n".join(lines)
    return [Document(page_content=text, metadata={"fuente": os.path.basename(path)})]


def load_json(path: str):
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    text = json.dumps(data, ensure_ascii=False, indent=2)
    return [Document(page_content=text, metadata={"fuente": os.path.basename(path)})]


def load_html(path: str):
    with open(path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    text = soup.get_text(separator="\n")
    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    return [Document(page_content=text, metadata={"fuente": os.path.basename(path)})]


# Mapa de extension -> funcion de carga
LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".xlsx": load_xlsx,
    ".pptx": load_pptx,
    ".md": load_markdown,
    ".csv": load_csv,
    ".json": load_json,
    ".html": load_html,
    ".htm": load_html,
}


def load_all_documents(data_dir: str):
    """Recorre data_dir, carga todos los archivos con formato soportado y
    devuelve la lista combinada de Document. Imprime un resumen por archivo."""
    all_docs = []
    files = sorted(glob.glob(os.path.join(data_dir, "*")))
    for path in files:
        if os.path.isdir(path):
            continue
        ext = os.path.splitext(path)[1].lower()
        loader_fn = LOADERS.get(ext)
        if not loader_fn:
            print(f"  Omitido (formato no soportado): {os.path.basename(path)}")
            continue
        try:
            docs = loader_fn(path)
            all_docs.extend(docs)
            print(f"  Cargado: {os.path.basename(path)} ({len(docs)} fragmento/s de origen)")
        except Exception as e:
            print(f"  ERROR al cargar {os.path.basename(path)}: {e}")
    return all_docs
